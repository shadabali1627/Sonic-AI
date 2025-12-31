import fitz  # PyMuPDF
import asyncio
from fastapi import UploadFile, HTTPException
from typing import Tuple
import docx
from pptx import Presentation
import io

class FileService:
    @staticmethod
    async def process_file(file: UploadFile) -> Tuple[str, str, bytes]:
        """
        Process uploaded file and return (file_type, content_or_path, raw_bytes).
        file_type: 'image', 'pdf', 'audio', 'text'
        """
        content = await file.read()
        filename = file.filename.lower()

        if filename.endswith(".pdf"):
            loop = asyncio.get_running_loop()
            text = await loop.run_in_executor(None, FileService._extract_text_from_pdf, content)
            return "pdf", text, content
        elif filename.endswith(".docx"):
            loop = asyncio.get_running_loop()
            text = await loop.run_in_executor(None, FileService._extract_text_from_docx, content)
            return "text", text, content
        elif filename.endswith(".pptx"):
            loop = asyncio.get_running_loop()
            text = await loop.run_in_executor(None, FileService._extract_text_from_pptx, content)
            return "text", text, content
        elif filename.endswith((".png", ".jpg", ".jpeg", ".webp")):
            return "image", "image_data", content
        elif filename.endswith((".wav", ".mp3", ".m4a")):
            return "audio", "audio_data", content
        elif filename.endswith((".txt", ".md", ".csv", ".json", ".py", ".js", ".ts", ".html", ".css")):
            # Basic text files
            try:
                return "text", content.decode('utf-8'), content
            except:
                return "binary", "binary_data", content
        else:
             try:
                return "text", content.decode('utf-8'), content
             except:
                return "binary", f"[Binary File: {filename}]", content

    @staticmethod
    def _extract_text_from_pdf(content: bytes) -> str:
        try:
            doc = fitz.open(stream=content, filetype="pdf")
            text = ""
            for page in doc:
                text += page.get_text()
            return text
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Error reading PDF: {str(e)}")

    @staticmethod
    def _extract_text_from_docx(content: bytes) -> str:
        try:
            doc = docx.Document(io.BytesIO(content))
            text = []
            for para in doc.paragraphs:
                text.append(para.text)
            return "\n".join(text)
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Error reading DOCX: {str(e)}")

    @staticmethod
    def _extract_text_from_pptx(content: bytes) -> str:
        try:
            prs = Presentation(io.BytesIO(content))
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Error reading PPTX: {str(e)}")
